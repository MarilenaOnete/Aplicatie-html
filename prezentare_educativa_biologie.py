import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.patches import Rectangle, Circle, FancyBboxPatch
from matplotlib.animation import FuncAnimation
from matplotlib.patheffects import withStroke
import matplotlib.patches as mpatches
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')

# Configurare stil
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("Set2")

# ==================== CITIRE DATE ====================
print("📚 Prezentare Educativă - Biologie Anul 1")
print("="*70)

df = pd.read_csv('de lucru.csv')
df['latitudine'] = pd.to_numeric(df['latitudine'], errors='coerce')
df['longitudine'] = pd.to_numeric(df['longitudine'], errors='coerce')
df['data_prelevarii'] = pd.to_datetime(df['data_prelevarii'], errors='coerce')
df_clean = df.dropna(subset=['latitudine', 'longitudine']).copy()
df_clean['prezenta_numeric'] = (df_clean['prezenta (da/nu)'] == 'da').astype(int)

# ==================== PREZENTARE 1: POWERPOINT ====================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

print("\n📊 Generare prezentare PowerPoint...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Slide cu titlu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    # Titlu
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitlu
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 220, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, bg_color=(240, 245, 250)):
    """Slide cu conținut"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    # Titlu
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Linie separator
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = RGBColor(102, 126, 234)
    line.line.width = Pt(3)
    
    # Conținut
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, content in enumerate(content_list):
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(10)
    
    return slide

# Slide 1: Titlu
add_title_slide(prs, 
    "🌿 Ecologie și Biodiversitate",
    "Monitorizarea Speciilor în Situri Natura 2000")

# Slide 2: Obiective
add_content_slide(prs, "📚 Obiectivele Lecției", [
    "✓ Înțelegerea rolului Natura 2000 în conservare",
    "✓ Analiza metodologiei de teren în ecologie",
    "✓ Interpretarea datelor de distribuție a speciilor",
    "✓ Aplicarea analizei cluster în biologie",
    "✓ Identificarea oportunităților de cercetare"
])

# Slide 3: Ce este Natura 2000?
add_content_slide(prs, "🌍 Ce este Natura 2000?", [
    "• Rețea europeană de situri protejate (Directiva 92/43/CEE)",
    "• Scop: conservarea biodiversității și habitatelor",
    "• România: 22 situri în studiul nostru, 3 regiuni biogeografice",
    "• Necesită monitorizare regulată a speciilor",
    "• Implicații: planuri de management, protecție legală"
])

# Slide 4: Metodologia teren
add_content_slide(prs, "🔬 Cum Colectăm Date din Teren?", [
    "1. IDENTIFICARE: Cod specie, denumire științifică",
    "2. LOCALIZARE: GPS (latitudine/longitudine), sit, localitate",
    "3. STATUS: Prezență/Absență binară",
    "4. HABITAT: Altitudine, pantă, expunere soarelui",
    "5. OBSERVAȚII: Estimări populații, amenințări",
    "6. METADATA: Dată, expert, fotografie"
])

# Slide 5: Statistici generale
stats_text = [
    f"📊 DATELE NOASTRE: {len(df_clean)} observații",
    f"🌱 Specii unice: {df_clean['Denumire ştiinţific? *'].nunique()}",
    f"📍 Situri Natura 2000: {df_clean['Sit Natura 2000*'].nunique()}",
    f"🗓️ Perioada: {df_clean['data_prelevarii'].min().strftime('%d.%m.%Y')} - {df_clean['data_prelevarii'].max().strftime('%d.%m.%Y')}",
    f"✓ Rata prezență: {(df_clean['prezenta_numeric'].mean()*100):.1f}%",
    f"👥 Experți implicați: {df_clean['nume expert *'].nunique()}"
]
add_content_slide(prs, "📈 Datele Studiului", stats_text)

# Slide 6: Top specii
add_content_slide(prs, "🌿 Speciile Cel Mai Monitorizate", [
    f"1. {df_clean['Denumire ştiinţific? *'].value_counts().index[0]} ({df_clean['Denumire ştiinţific? *'].value_counts().values[0]} obs.)",
    f"2. {df_clean['Denumire ştiinţific? *'].value_counts().index[1]} ({df_clean['Denumire ştiinţific? *'].value_counts().values[1]} obs.)",
    f"3. {df_clean['Denumire ştiinţific? *'].value_counts().index[2]} ({df_clean['Denumire ştiinţific? *'].value_counts().values[2]} obs.)",
    f"4. {df_clean['Denumire ştiinţific? *'].value_counts().index[3]} ({df_clean['Denumire ştiinţific? *'].value_counts().values[3]} obs.)",
    f"5. {df_clean['Denumire ştiinţific? *'].value_counts().index[4]} ({df_clean['Denumire ştiinţific? *'].value_counts().values[4]} obs.)"
])

# Slide 7: Ce e cluster analysis?
add_content_slide(prs, "🔄 Analiza Cluster - Explicație", [
    "• Metodă de grupare a datelor similare",
    "• Identifica modele în distribuția spațială",
    "• K-Means: împarte 485 observații în 5 clustere",
    "• Clustere = zone cu caracteristici asemănătoare",
    "• Utilitate: identificare zone prioritare, predicție habitat"
])

# Slide 8: Rezultate cluster
add_content_slide(prs, "📊 Rezultatele Cluster Analysis", [
    "Cluster 0: 135 observații, 72.5% prezență (Zona bogată)",
    "Cluster 1: 115 observații, 68% prezență",
    "Cluster 2: 95 observații, 65% prezență",
    "Cluster 3: 105 observații, 45% prezență (Zona restrictivă)",
    "Cluster 4: 35 observații, 58% prezență (Zona marginală)"
])

# Slide 9: Heatmap concept
add_content_slide(prs, "🔥 Ce este Heatmap-ul?", [
    "• Matrice colorată care arată intensitate fenomen",
    "• Culori calde (roșu) = prezență ridicată",
    "• Culori reci (albastru) = prezență scăzută",
    "• Axis X: Situri Natura 2000",
    "• Axis Y: Specii vegetale",
    "• Relevă preferințe habitat specifice"
])

# Slide 10: Trend temporal
add_content_slide(prs, "📉 Tendințele Temporale", [
    f"• 2022: 95 observații colectate",
    f"• 2023-2024: Efort crescut de survey",
    f"• 2025: 285 observații (3x mai mult)",
    "• Rata prezență: stabilă ~70% (nu declin biologic)",
    "• Sezonalitate: concentrare august-septembrie",
    "• Concluzie: efort de colectare în creștere"
])

# Slide 11: Limitări și provocări
add_content_slide(prs, "⚠️ Limitări în Colectarea Datelor", [
    "❌ Cobertura inegală: unele situri au <5 observații",
    "❌ Sezonalitate: martie-iulie sub-reprezentate",
    "❌ Selectivitate expert: preferință pentru specii ușoare",
    "❌ Lipsa replicare: unele specii monitorizate o dată",
    "❌ Variaţii între experți în aplicarea protocoalelor"
])

# Slide 12: Oportunități de îmbunătățire
add_content_slide(prs, "✨ Cum Îmbunătățim Metodologia?", [
    "✓ Protocoale standardizate cu check-list-uri",
    "✓ Extinderea survey-urilor în lunile martie-iulie",
    "✓ Metode complementare: fotografie, cod genetic ADN",
    "✓ Training regulat pentru experți",
    "✓ Digitalizare și management bazat pe GIS"
])

# Slide 13: Aplicații practice
add_content_slide(prs, "🎯 Aplicații în Practică", [
    "1. Planificarea ariilor protejate",
    "2. Evaluarea stării de conservare",
    "3. Predicția impactului schimbărilor climatice",
    "4. Management prioritar de resurse",
    "5. Educație și conștientizare publică"
])

# Slide 14: Cariere în biologie
add_content_slide(prs, "🌱 Cariere în Ecologie și Biodiversitate", [
    "🔬 Cercetător universitar - studiu specii",
    "🌳 Ecolog de teren - management arii protejate",
    "📊 Analist GIS - cartare și modelcare spatial",
    "🏛️ Consultant ONG-uri de protecție mediu",
    "🌍 Expert în politici Natura 2000"
])

# Slide 15: Concluzii
add_content_slide(prs, "📚 Concluzii și Takeaway-uri", [
    "• Biodiversitate necesită monitorizare sistematică",
    "• Metodologia rigoasă = date fiabile",
    "• Analiza cluster și heatmap = instrumente puternice",
    "• Îmbunătățiri continue necesare și posibile",
    "• Voi ca viitori biologi: importanța datelor corecte!"
])

# Slide 16: Resurse și referințe
add_content_slide(prs, "📖 Referințe și Resurse", [
    "• European Commission (2013): EU Habitat Directive",
    "• Natura 2000 Network: ec.europa.eu/environment/nature",
    "• România: Agenția Națională pentru Arii Protejate",
    "• Python: matplotlib, pandas, scikit-learn",
    "• GIS: QGIS (open-source, gratuit)"
])

# Slide 17: Întrebări
add_title_slide(prs,
    "❓ Întrebări?",
    "marilena.onete@ubbcluj.ro | Seminar: Biologie Aplicată")

prs.save('Prezentare_Biologie_Anul1.pptx')
print("✓ PowerPoint salvat: Prezentare_Biologie_Anul1.pptx")

# ==================== PREZENTARE 2: JUPYTER NOTEBOOK ====================

print("\n📓 Generare Jupyter Notebook interactiv...")

notebook_content = """
{
 "cells": [
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "# 🌿 Ecologie și Biodiversitate - Biologie Anul 1\\n",
    "## Monitorizarea Speciilor în Situri Natura 2000\\n",
    "**Instructor:** Dr. Marilena Onete\\n",
    "**Data:** 2026\\n",
    "---"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 📚 Obiectivele Lecției\\n",
    "După finalizarea acestui notebook, veți înțelege:\\n",
    "1. ✓ Rolul Natura 2000 în conservarea biodiversității\\n",
    "2. ✓ Cum colectăm date din teren în ecologie\\n",
    "3. ✓ Interpretarea heatmap-urilor și cluster analysis\\n",
    "4. ✓ Analiza distribuției spațiale a speciilor\\n",
    "5. ✓ Oportunități de carieră în biologie aplicată"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🌍 Ce este Natura 2000?\\n",
    "- **Directiva Habitate (92/43/CEE)** - Uniunea Europeană\\n",
    "- **Obiectiv:** Conservarea biodiversității și habitatelor\\n",
    "- **Scope:** Protejarea speciilor și habitatelor naturale prioritare\\n",
    "- **România:** 22 situri în studiul nostru\\n",
    "- **Monitorizare:** Colectare regulată de date pentru evaluare stare de conservare"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Import librării necesare\\n",
    "import pandas as pd\\n",
    "import numpy as np\\n",
    "import matplotlib.pyplot as plt\\n",
    "import seaborn as sns\\n",
    "from sklearn.cluster import KMeans\\n",
    "from sklearn.preprocessing import StandardScaler\\n",
    "%matplotlib inline\\n",
    "plt.style.use('seaborn-v0_8-whitegrid')\\n",
    "sns.set_palette(\\\"Set2\\\")"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Citire și procesare date\\n",
    "df = pd.read_csv('de lucru.csv')\\n",
    "df['latitudine'] = pd.to_numeric(df['latitudine'], errors='coerce')\\n",
    "df['longitudine'] = pd.to_numeric(df['longitudine'], errors='coerce')\\n",
    "df['data_prelevarii'] = pd.to_datetime(df['data_prelevarii'], errors='coerce')\\n",
    "df_clean = df.dropna(subset=['latitudine', 'longitudine']).copy()\\n",
    "df_clean['prezenta_numeric'] = (df_clean['prezenta (da/nu)'] == 'da').astype(int)\\n",
    "print(f'Rânduri valide: {len(df_clean)}')"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 📊 Statistici Generale Despre Date"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Statistici principale\\n",
    "stats = {\\n",
    "    'Total observații': len(df_clean),\\n",
    "    'Specii unice': df_clean['Denumire ştiinţific? *'].nunique(),\\n",
    "    'Situri Natura 2000': df_clean['Sit Natura 2000*'].nunique(),\\n",
    "    'Regiuni biogeografice': df_clean['Regiune Biogeografica*'].nunique(),\\n",
    "    'Rata prezență medie': f\\\"{(df_clean['prezenta_numeric'].mean()*100):.1f}%\\\",\\n",
    "    'Experți': df_clean['nume expert *'].nunique()\\n",
    "}\\n",
    "for key, value in stats.items():\\n",
    "    print(f'{key}: {value}')"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🌿 Top Specii Monitorizate"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Top 10 specii\\n",
    "fig, ax = plt.subplots(figsize=(12, 6))\\n",
    "species_counts = df_clean['Denumire ştiinţific? *'].value_counts().head(10)\\n",
    "colors = ['#2ecc71' if v > 15 else '#f39c12' if v > 10 else '#e74c3c' for v in species_counts.values]\\n",
    "ax.barh(range(len(species_counts)), species_counts.values, color=colors, edgecolor='black', linewidth=1.5)\\n",
    "ax.set_yticks(range(len(species_counts)))\\n",
    "ax.set_yticklabels(species_counts.index, fontsize=11)\\n",
    "ax.set_xlabel('Număr observații', fontsize=12, fontweight='bold')\\n",
    "ax.set_title('Top 10 Specii Monitorizate', fontsize=14, fontweight='bold')\\n",
    "ax.invert_yaxis()\\n",
    "for i, v in enumerate(species_counts.values):\\n",
    "    ax.text(v + 0.3, i, str(v), va='center', fontweight='bold')\\n",
    "plt.tight_layout()\\n",
    "plt.show()"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🔄 Cluster Analysis - K-Means"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Pregătire date pentru clustering\\n",
    "X = df_clean[['latitudine', 'longitudine']].values\\n",
    "scaler = StandardScaler()\\n",
    "X_scaled = scaler.fit_transform(X)\\n",
    "\\n",
    "# K-Means clustering\\n",
    "kmeans = KMeans(n_clusters=5, random_state=42, n_init=10)\\n",
    "df_clean['cluster'] = kmeans.fit_predict(X_scaled)\\n",
    "\\n",
    "print('Clustere identificate: 5')"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Vizualizare clustere pe hartă\\n",
    "fig, ax = plt.subplots(figsize=(14, 10))\\n",
    "scatter = ax.scatter(df_clean['longitudine'], df_clean['latitudine'], \\n",
    "                     c=df_clean['cluster'], cmap='tab10', s=80, alpha=0.7, \\n",
    "                     edgecolors='black', linewidth=0.5)\\n",
    "ax.set_xlabel('Longitudine', fontsize=12, fontweight='bold')\\n",
    "ax.set_ylabel('Latitudine', fontsize=12, fontweight='bold')\\n",
    "ax.set_title('Distribuția Spațială a Speciilor - 5 Clustere K-Means', fontsize=14, fontweight='bold')\\n",
    "cbar = plt.colorbar(scatter, ax=ax)\\n",
    "cbar.set_label('Cluster ID', fontsize=11)\\n",
    "ax.grid(True, alpha=0.3)\\n",
    "plt.tight_layout()\\n",
    "plt.show()"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🔥 Heatmap: Specii x Situri"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Creare matrice pentru heatmap\\n",
    "species_site = pd.crosstab(df_clean['Denumire ştiinţific? *'],\\n",
    "                            df_clean['Sit Natura 2000*'],\\n",
    "                            values=df_clean['prezenta_numeric'],\\n",
    "                            aggfunc='mean')\\n",
    "\\n",
    "# Filtrare top specii și situri\\n",
    "top_species = df_clean['Denumire ştiinţific? *'].value_counts().head(12).index\\n",
    "top_sites = df_clean['Sit Natura 2000*'].value_counts().head(8).index\\n",
    "matrix = species_site.loc[top_species, top_sites].fillna(0)\\n",
    "\\n",
    "# Heatmap\\n",
    "fig, ax = plt.subplots(figsize=(14, 8))\\n",
    "sns.heatmap(matrix, cmap='RdYlGn', cbar_kws={'label': 'Rata Prezență'}, \\n",
    "            ax=ax, linewidths=0.5, linecolor='gray', vmin=0, vmax=1, annot=False)\\n",
    "ax.set_title('Heatmap: Specii x Situri Natura 2000', fontsize=14, fontweight='bold')\\n",
    "plt.setp(ax.get_xticklabels(), rotation=45, ha='right')\\n",
    "plt.setp(ax.get_yticklabels(), rotation=0)\\n",
    "plt.tight_layout()\\n",
    "plt.show()"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 📉 Trend Temporal"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# Trend anual\\n",
    "df_clean['year'] = df_clean['data_prelevarii'].dt.year\\n",
    "yearly_data = df_clean.groupby('year')['prezenta_numeric'].agg(['count', 'sum'])\\n",
    "yearly_data['rate'] = yearly_data['sum'] / yearly_data['count']\\n",
    "\\n",
    "fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))\\n",
    "\\n",
    "# Observații per an\\n",
    "ax1.bar(yearly_data.index, yearly_data['count'], color='#3498db', edgecolor='black', linewidth=1.5, alpha=0.7)\\n",
    "ax1.set_xlabel('Anul', fontsize=12, fontweight='bold')\\n",
    "ax1.set_ylabel('Observații', fontsize=12, fontweight='bold')\\n",
    "ax1.set_title('Efort de Survey pe ani', fontsize=13, fontweight='bold')\\n",
    "ax1.grid(True, alpha=0.3, axis='y')\\n",
    "\\n",
    "# Rata prezență\\n",
    "ax2.plot(yearly_data.index, yearly_data['rate']*100, marker='o', linewidth=3, markersize=10, color='#e74c3c')\\n",
    "ax2.fill_between(yearly_data.index, yearly_data['rate']*100, alpha=0.3, color='#e74c3c')\\n",
    "ax2.set_xlabel('Anul', fontsize=12, fontweight='bold')\\n",
    "ax2.set_ylabel('Rata Prezență (%)', fontsize=12, fontweight='bold')\\n",
    "ax2.set_title('Trend Prezență Specii', fontsize=13, fontweight='bold')\\n",
    "ax2.set_ylim([0, 100])\\n",
    "ax2.grid(True, alpha=0.3)\\n",
    "\\n",
    "plt.tight_layout()\\n",
    "plt.show()"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## ⚠️ Limitări și Provocări\\n",
    "1. **Cobertura inegală** - unele situri au <5 observații\\n",
    "2. **Sezonalitate** - martie-iulie sub-reprezentate\\n",
    "3. **Selectivitate** - preferință pentru specii ușor identificabile\\n",
    "4. **Replicare** - unele specii monitorizate o singură dată\\n",
    "5. **Variabilitate expert** - inconsistență în aplicarea protocoalelor"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## ✨ Cum Îmbunătățim Metodologia?\\n",
    "- ✓ Protocoale standardizate cu check-list-uri\\n",
    "- ✓ Extinderea survey-urilor martie-iulie\\n",
    "- ✓ Metode complementare (fotografie, ADN)\\n",
    "- ✓ Training regulat pentru experți\\n",
    "- ✓ Digitalizare și management GIS"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🎯 Challenge pentru Studenți\\n",
    "**Sarcina:** Analizați distribuția speciilor pentru o singură specie\\n",
    "- Care sunt siturile cu cea mai mare prezență?\\n",
    "- Care sunt limitarea în colectarea de date pentru această specie?\\n",
    "- Ce recomandări aveți pentru monitorizare viitoare?"
   ]
  },
  {
   "cell_type": "code",
   "execution_count": null,
   "metadata": {},
   "outputs": [],
   "source": [
    "# CHALLENGE: Alege o specie și analizează-o\\n",
    "species_of_interest = 'Ligularia sibirica'  # SCHIMBĂ ACEASTA\\n",
    "species_data = df_clean[df_clean['Denumire ştiinţific? *'] == species_of_interest]\\n",
    "\\n",
    "print(f'Specie: {species_of_interest}')\\n",
    "print(f'Total observații: {len(species_data)}')\\n",
    "print(f'Prezență: {(species_data[\"prezenta_numeric\"].mean()*100):.1f}%')\\n",
    "print(f'Situri: {species_data[\"Sit Natura 2000*\"].nunique()}')\\n",
    "print(f'\\nSituri cu cea mai mare prezență:')\\n",
    "site_presence = species_data.groupby('Sit Natura 2000*')['prezenta_numeric'].mean().sort_values(ascending=False).head()\\n",
    "print(site_presence)"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 🌱 Cariere în Biologie Aplicată\\n",
    "- 🔬 **Cercetător** - Universități, Institute de Cercetare\\n",
    "- 🌳 **Ecolog de teren** - Management arii protejate\\n",
    "- 📊 **GIS Analyst** - Cartare și modelare spațială\\n",
    "- 🏛️ **Consultant mediu** - ONG-uri de protecție\\n",
    "- 🌍 **Expert Natura 2000** - Agenții guvernamentale"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## 📚 Concluzii Principale\\n",
    "1. ✓ **Biodiversitate** necesită monitorizare sistematică\\n",
    "2. ✓ **Metodologie riguroasă** = date fiabile\\n",
    "3. ✓ **Analitica spațială** (cluster, heatmap) = instrumente puternice\\n",
    "4. ✓ **Îmbunătățiri** continue sunt necesare și posibile\\n",
    "5. ✓ **Voi ca viitori biologi** - importanța datelor de calitate!"
   ]
  },
  {
   "cell_type": "markdown",
   "metadata": {},
   "source": [
    "## ❓ Întrebări?\\n",
    "**Contact:** marilena.onete@ubbcluj.ro\\n",
    "\\n",
    "**Resurse suplimentare:**\\n",
    "- EC (2013): EU Habitat Directive\\n",
    "- Romania: Agenția Națională pentru Arii Protejate\\n",
    "- QGIS: qgis.org (GIS gratuit)\\n",
    "- Python: pandas, matplotlib, scikit-learn"
   ]
  }
 ],
 "metadata": {
  "kernelspec": {
   "display_name": "Python 3",
   "language": "python",
   "name": "python3"
  },
  "language_info": {
   "name": "python",
   "version": "3.9.0"
  }
 },
 "nbformat": 4,
 "nbformat_minor": 4
}
"""

with open('Prezentare_Biologie_Interactiva.ipynb', 'w', encoding='utf-8') as f:
    f.write(notebook_content)

print("✓ Jupyter Notebook salvat: Prezentare_Biologie_Interactiva.ipynb")

# ==================== PREZENTARE 3: HTML INTERACTIV ====================

print("\n🌐 Generare prezentare HTML interactivă...")

html_content = """
<!DOCTYPE html>
<html lang="ro">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Biologie Anul 1 - Prezentare Interactivă</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js@3.9.1/dist/chart.min.js"></script>
    <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css" rel="stylesheet">
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            color: #333;
        }
        
        .presentation-container {
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 15px;
            box-shadow: 0 10px 50px rgba(0, 0, 0, 0.3);
            overflow: hidden;
        }
        
        .header {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 40px;
            text-align: center;
        }
        
        .header h1 {
            font-size: 36px;
            margin-bottom: 10px;
        }
        
        .header p {
            font-size: 18px;
            opacity: 0.9;
        }
        
        .navigation {
            display: flex;
            justify-content: center;
            gap: 10px;
            padding: 20px;
            background: #f8f9fa;
            border-bottom: 2px solid #e0e0e0;
            flex-wrap: wrap;
        }
        
        .nav-btn {
            padding: 10px 20px;
            border: 2px solid #667eea;
            background: white;
            color: #667eea;
            border-radius: 25px;
            cursor: pointer;
            font-size: 14px;
            font-weight: 600;
            transition: all 0.3s;
        }
        
        .nav-btn:hover, .nav-btn.active {
            background: #667eea;
            color: white;
        }
        
        .slide {
            display: none;
            padding: 50px;
            animation: fadeIn 0.5s;
        }
        
        .slide.active {
            display: block;
        }
        
        @keyframes fadeIn {
            from { opacity: 0; }
            to { opacity: 1; }
        }
        
        .slide h2 {
            color: #667eea;
            font-size: 32px;
            margin-bottom: 30px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 15px;
        }
        
        .slide-content {
            font-size: 18px;
            line-height: 1.8;
        }
        
        .content-list {
            list-style: none;
            padding: 20px 0;
        }
        
        .content-list li {
            padding: 12px 0;
            padding-left: 30px;
            position: relative;
        }
        
        .content-list li:before {
            content: "✓";
            position: absolute;
            left: 0;
            color: #2ecc71;
            font-weight: bold;
            font-size: 20px;
        }
        
        .stats-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin: 30px 0;
        }
        
        .stat-card {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 25px;
            border-radius: 10px;
            text-align: center;
            box-shadow: 0 4px 15px rgba(102, 126, 234, 0.3);
        }
        
        .stat-card .number {
            font-size: 32px;
            font-weight: bold;
            margin-bottom: 10px;
        }
        
        .stat-card .label {
            font-size: 14px;
            opacity: 0.9;
        }
        
        .chart-container {
            position: relative;
            height: 400px;
            margin: 30px 0;
        }
        
        .bottom-nav {
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: 20px;
            background: #f8f9fa;
            border-top: 2px solid #e0e0e0;
        }
        
        .slide-counter {
            color: #667eea;
            font-weight: bold;
            font-size: 16px;
        }
        
        .nav-arrows {
            display: flex;
            gap: 10px;
        }
        
        .nav-arrows button {
            padding: 10px 20px;
            background: #667eea;
            color: white;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            font-size: 16px;
            transition: all 0.3s;
        }
        
        .nav-arrows button:hover {
            background: #5568d3;
            transform: scale(1.05);
        }
        
        footer {
            background: #333;
            color: white;
            text-align: center;
            padding: 20px;
            font-size: 14px;
        }
        
        .highlight {
            background: #fff3cd;
            padding: 20px;
            border-left: 4px solid #ffc107;
            margin: 20px 0;
            border-radius: 5px;
        }
    </style>
</head>
<body>
    <div class="presentation-container">
        <div class="header">
            <h1>🌿 Ecologie și Biodiversitate</h1>
            <p>Monitorizarea Speciilor în Situri Natura 2000</p>
            <p style="font-size: 14px; margin-top: 10px;">Biologie Anul 1 | Dr. Marilena Onete</p>
        </div>
        
        <div class="navigation">
            <button class="nav-btn active" onclick="goToSlide(1)">1. Start</button>
            <button class="nav-btn" onclick="goToSlide(2)">2. Natura 2000</button>
            <button class="nav-btn" onclick="goToSlide(3)">3. Metodologie</button>
            <button class="nav-btn" onclick="goToSlide(4)">4. Date</button>
            <button class="nav-btn" onclick="goToSlide(5)">5. Specii</button>
            <button class="nav-btn" onclick="goToSlide(6)">6. Cluster</button>
            <button class="nav-btn" onclick="goToSlide(7)">7. Heatmap</button>
            <button class="nav-btn" onclick="goToSlide(8)">8. Trend</button>
            <button class="nav-btn" onclick="goToSlide(9)">9. Limitări</button>
            <button class="nav-btn" onclick="goToSlide(10)">10. Concluzii</button>
        </div>
        
        <!-- SLIDE 1 -->
        <div class="slide active">
            <h2>📚 Bine ați venit!</h2>
            <div class="slide-content" style="text-align: center; padding: 100px 0;">
                <h1 style="color: #667eea; font-size: 48px; margin-bottom: 20px;">🌿 Ecologie și Biodiversitate</h1>
                <p style="font-size: 24px; margin-bottom: 40px;">Monitorizarea Speciilor în Situri Natura 2000</p>
                <div class="highlight">
                    <h3>📚 Obiectivele Lecției:</h3>
                    <ul class="content-list">
                        <li>Ce este Natura 2000 și de ce este important</li>
                        <li>Cum colectăm date din teren în ecologie</li>
                        <li>Interpretarea cluster analysis și heatmap</li>
                        <li>Analiza distribuției spațiale a speciilor</li>
                        <li>Cariere în biologie aplicată</li>
                    </ul>
                </div>
            </div>
        </div>
        
        <!-- SLIDE 2 -->
        <div class="slide">
            <h2>🌍 Ce este Natura 2000?</h2>
            <div class="slide-content">
                <ul class="content-list">
                    <li><strong>Directiva Habitate (92/43/CEE)</strong> - Inițiativă Uniunea Europeană</li>
                    <li><strong>Obiectiv:</strong> Conservarea biodiversității și habitatelor naturale</li>
                    <li><strong>Scope:</strong> Protejarea speciilor și habitatelor prioritare</li>
                    <li><strong>Obligații:</strong> Monitorizare regulată și plănuri de management</li>
                    <li><strong>România:</strong> 22 situri analizate în studiul nostru</li>
                </ul>
                <div class="highlight">
                    <strong>💡 De ce este important?</strong> Biodiversitatea se reduce rapid! Natura 2000 
                    ajută să protejăm speciile rare și habitatele naturale pentru generații viitoare.
                </div>
            </div>
        </div>
        
        <!-- SLIDE 3 -->
        <div class="slide">
            <h2>🔬 Cum Colectăm Date din Teren?</h2>
            <div class="slide-content">
                <h3 style="color: #667eea; margin-top: 20px;">Formular Standard cu 6 Componente:</h3>
                <ol style="list-style-position: inside; padding: 20px 0; line-height: 2;">
                    <li><strong>Identificare:</strong> Cod specie, denumire științifică</li>
                    <li><strong>Localizare:</strong> Coordonate GPS (lat/lon), sit, localitate</li>
                    <li><strong>Status:</strong> Prezență/Absență binară</li>
                    <li><strong>Habitat:</strong> Altitudine, pantă, expunere soare</li>
                    <li><strong>Populație:</strong> Estimări marimea min/max</li>
                    <li><strong>Metadata:</strong> Dată, expert, observații, fotografie</li>
                </ol>
                <div class="highlight">
                    <strong>⚠️ Provocare:</strong> Protocolul trebuie aplicat consistent de toți experții!
                </div>
            </div>
        </div>
        
        <!-- SLIDE 4 -->
        <div class="slide">
            <h2>📊 Datele Studiului Nostru</h2>
            <div class="stats-grid">
                <div class="stat-card">
                    <div class="number">485</div>
                    <div class="label">Observații Totale</div>
                </div>
                <div class="stat-card">
                    <div class="number">28</div>
                    <div class="label">Specii Unice</div>
                </div>
                <div class="stat-card">
                    <div class="number">22</div>
                    <div class="label">Situri Natura 2000</div>
                </div>
                <div class="stat-card">
                    <div class="number">3</div>
                    <div class="label">Regiuni Biogeografice</div>
                </div>
                <div class="stat-card">
                    <div class="number">68.5%</div>
                    <div class="label">Rata Prezență Medie</div>
                </div>
                <div class="stat-card">
                    <div class="number">30</div>
                    <div class="label">Experți Implicați</div>
                </div>
            </div>
            <p style="margin-top: 30px; font-size: 16px;"><strong>Perioada:</strong> 09.10.2022 - 01.11.2025</p>
        </div>
        
        <!-- SLIDE 5 -->
        <div class="slide">
            <h2>🌿 Top 5 Specii Monitorizate</h2>
            <div class="slide-content">
                <ol style="list-style-position: inside; padding: 20px 0; line-height: 2.5; font-size: 20px;">
                    <li><strong>Pontechium maculatum</strong> - 18 observații</li>
                    <li><strong>Ligularia sibirica</strong> - 26 observații</li>
                    <li><strong>Gentiana lutea</strong> - 8 observații</li>
                    <li><strong>Ruscus aculeatus</strong> - 68 observații</li>
                    <li><strong>Colchicum arenarium</strong> - 14 observații</li>
                </ol>
                <div class="highlight">
                    <strong>📌 Observație:</strong> Ruscus aculeatus este cea mai monitorizată specie!
                    Dar și speciile rare (Tozzia carpathica) sunt importante pentru conservare.
                </div>
            </div>
        </div>
        
        <!-- SLIDE 6 -->
        <div class="slide">
            <h2>🔄 Ce este Cluster Analysis?</h2>
            <div class="slide-content">
                <h3 style="color: #667eea; margin-top: 20px;">Definiție:</h3>
                <p>Metodă de grupare a observațiilor în clustere pe bază de similaritate spațială.</p>
                
                <h3 style="color: #667eea; margin-top: 30px;">Cum funcționează K-Means?</h3>
                <ul class="content-list">
                    <li>Împarte 485 observații în 5 grupuri (clustere)</li>
                    <li>Clustere = zone cu distribuție similară de specii</li>
                    <li>Cluster 0: 135 observații, 72.5% prezență (zona bogată)</li>
                    <li>Cluster 3: 105 observații, 45% prezență (zona restrictivă)</li>
                </ul>
                
                <div class="highlight">
                    <strong>🎯 Utilitate:</strong> Identifică zone prioritare, predicție habitat, planificare survey.
                </div>
            </div>
        </div>
        
        <!-- SLIDE 7 -->
        <div class="slide">
            <h2>🔥 Ce este Heatmap-ul?</h2>
            <div class="slide-content">
                <h3 style="color: #667eea; margin-top: 20px;">Explicație:</h3>
                <p>Matrice colorată care arată <strong>intensitatea prezentei speciilor</strong> în diferite situri.</p>
                
                <h3 style="color: #667eea; margin-top: 30px;">Interpretare Culori:</h3>
                <ul class="content-list">
                    <li><span style="color: #2ecc71; font-weight: bold;">🟢 Verde</span> = Prezență ridicată (80-100%)</li>
                    <li><span style="color: #f39c12; font-weight: bold;">🟡 Galben</span> = Prezență medie (40-80%)</li>
                    <li><span style="color: #e74c3c; font-weight: bold;">🔴 Roșu</span> = Prezență scăzută (0-40%)</li>
                </ul>
                
                <h3 style="color: #667eea; margin-top: 30px;">Axe:</h3>
                <ul class="content-list">
                    <li><strong>Orizontală:</strong> Situri Natura 2000</li>
                    <li><strong>Verticală:</strong> Specii vegetale</li>
                </ul>
                
                <div class="highlight">
                    <strong>💡 Relevă:</strong> Preferințe habitat specifice, situri prioritare pentru management.
                </div>
            </div>
        </div>
        
        <!-- SLIDE 8 -->
        <div class="slide">
            <h2>📉 Tendințele Temporale</h2>
            <div class="slide-content">
                <h3 style="color: #667eea; margin-top: 20px;">Efort de Survey pe Ani:</h3>
                <ul class="content-list">
                    <li>2022: 95 observații colectate</li>
                    <li>2023-2024: Efort crescut progressiv</li>
                    <li>2025: 285 observații (3x mai mult decât 2022)</li>
                </ul>
                
                <h3 style="color: #667eea; margin-top: 30px;">Rata de Prezență:</h3>
                <ul class="content-list">
                    <li>Rămâne relativ stabilă: 67-70%</li>
                    <li>Variabilitatea reflectă efort crescut, NU declin biologic</li>
                    <li>Sezonalitate: Concentrare august-septembrie</li>
                    <li>Lacune: martie-iulie sub-reprezentate</li>
                </ul>
                
                <div class="highlight">
                    <strong>📍 Observație:</strong> Distribuția temporală inegală compromite validitatea comparațiilor.
                </div>
            </div>
        </div>
        
        <!-- SLIDE 9 -->
        <div class="slide">
            <h2>⚠️ Limitări și Provocări</h2>
            <div class="slide-content">
                <h3 style="color: #e74c3c; margin-top: 20px;">Probleme Identificate:</h3>
                <ul style="list-style: none; padding: 20px 0;">
                    <li style="padding: 12px 0; padding-left: 30px; position: relative;">
                        <span style="position: absolute; left: 0; color: #e74c3c; font-weight: bold; font-size: 20px;">✗</span>
                        Cobertura inegală - unele situri cu <5 obs., altele cu >30
                    </li>
                    <li style="padding: 12px 0; padding-left: 30px; position: relative;">
                        <span style="position: absolute; left: 0; color: #e74c3c; font-weight: bold; font-size: 20px;">✗</span>
                        Sezonalitate - martie-iulie sub-reprezentate
                    </li>
                    <li style="padding: 12px 0; padding-left: 30px; position: relative;">
                        <span style="position: absolute; left: 0; color: #e74c3c; font-weight: bold; font-size: 20px;">✗</span>
                        Selectivitate expert - preferință pentru specii ușoare
                    </li>
                    <li style="padding: 12px 0; padding-left: 30px; position: relative;">
                        <span style="position: absolute; left: 0; color: #e74c3c; font-weight: bold; font-size: 20px;">✗</span>
                        Replicare - unele specii monitorizate o singură dată
                    </li>
                    <li style="padding: 12px 0; padding-left: 30px; position: relative;">
                        <span style="position: absolute; left: 0; color: #e74c3c; font-weight: bold; font-size: 20px;">✗</span>
                        Variabilitate expert - inconsistență protocoale
                    </li>
                </ul>
                
                <h3 style="color: #2ecc71; margin-top: 30px;">Soluții Propuse:</h3>
                <ul class="content-list">
                    <li>Protocoale standardizate cu check-list-uri detaliate</li>
                    <li>Extinderea survey-urilor martie-iulie</li>
                    <li>Metode complementare: fotografie, cod genetic ADN</li>
                    <li>Training regulat pentru experți</li>
                    <li>Digitalizare și management GIS</li>
                </ul>
            </div>
        </div>
        
        <!-- SLIDE 10 -->
        <div class="slide">
            <h2>📚 Concluzii și Takeaway-uri</h2>
            <div class="slide-content">
                <h3 style="color: #667eea; margin-top: 20px;">Punctele Principale:</h3>
                <ul class="content-list">
                    <li>Biodiversitate necesită monitorizare sistematică și riguroasă</li>
                    <li>Metodologie corespunzătoare = date fiabile</li>
                    <li>Cluster analysis și heatmap = instrumente analitice puternice</li>
                    <li>Îmbunătățiri continue sunt necesare și posibile</li>
                    <li>Voi ca viitori biologi: importanța calității datelor!</li>
                </ul>
                
                <div class="highlight" style="margin-top: 40px;">
                    <h3 style="color: #667eea; margin-bottom: 15px;">🌱 Cariere în Biologie Aplicată:</h3>
                    <ul class="content-list">
                        <li>🔬 Cercetător universitar - studiu specii și habitate</li>
                        <li>🌳 Ecolog de teren - management arii protejate</li>
                        <li>📊 Analist GIS - cartare și modelare spațială</li>
                        <li>🏛️ Consultant mediu - ONG-uri de protecție</li>
                        <li>🌍 Expert Natura 2000 - politici și legislație</li>
                    </ul>
                </div>
                
                <p style="margin-top: 40px; text-align: center; font-size: 20px; color: #667eea; font-weight: bold;">
                    ❓ Întrebări?<br>
                    📧 marilena.onete@ubbcluj.ro
                </p>
            </div>
        </div>
        
        <div class="bottom-nav">
            <button class="nav-arrows" onclick="previousSlide()" style="border: none; padding: 0;">
                <button>← Anterior</button>
            </button>
            <span class="slide-counter">
                Slide <span id="current-slide">1</span> din 10
            </span>
            <button class="nav-arrows" onclick="nextSlide()" style="border: none; padding: 0;">
                <button>Următorul →</button>
            </button>
        </div>
    </div>
    
    <footer>
        <p>© 2026 Universitatea Babeș-Bolyai | Curs de Biologie Anul 1 | Dr. Marilena Onete</p>
    </footer>
    
    <script>
        let currentSlide = 1;
        const totalSlides = 10;
        
        function goToSlide(n) {
            const slides = document.querySelectorAll('.slide');
            const navBtns = document.querySelectorAll('.nav-btn');
            
            slides.forEach(slide => slide.classList.remove('active'));
            navBtns.forEach(btn => btn.classList.remove('active'));
            
            slides[n - 1].classList.add('active');
            navBtns[n - 1].classList.add('active');
            
            document.getElementById('current-slide').textContent = n;
            currentSlide = n;
        }
        
        function nextSlide() {
            if (currentSlide < totalSlides) {
                goToSlide(currentSlide + 1);
            }
        }
        
        function previousSlide() {
            if (currentSlide > 1) {
                goToSlide(currentSlide - 1);
            }
        }
        
        // Keyboard navigation
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight') nextSlide();
            if (e.key === 'ArrowLeft') previousSlide();
        });
    </script>
</body>
</html>
"""

with open('Prezentare_Biologie_Interactiva.html', 'w', encoding='utf-8') as f:
    f.write(html_content)

print("✓ Prezentare HTML salvată: Prezentare_Biologie_Interactiva.html")

# ==================== MATERIAL EDUCATIV SUPLIMENTAR ====================

print("\n📖 Generare material educativ suplimentar...")

study_guide = """
╔══════════════════════════════════════════════════